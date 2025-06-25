import os
import subprocess
import json
import re
from pptx import Presentation

# 1) SET YOUR OPENAI API KEY
OPENAI_API_KEY = "INSERT_YOUR_OPENAI_API_KEY_HERE"

# 2) FUNCTION TO EXTRACT SLIDE TEXT
def extract_slides_text(pptx_path):
    prs = Presentation(pptx_path)
    slides_text = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        slides_text.append("\n".join(slide_text))
    return slides_text

# 3) CHUNKING FUNCTION
def chunk_list(lst, chunk_size):
    """Yield successive chunk_size-sized chunks from lst."""
    for i in range(0, len(lst), chunk_size):
        yield lst[i:i + chunk_size]

# 4) FUNCTION TO GET SPEAKER NOTES USING CURL WITH BETTER ERROR HANDLING
def get_speaker_notes(slides_text_chunk):
    # Construct the prompt for the API
    # I have included a generic prompt template, but I've found it works better when you give it the general topic of the slides (e.g. "cognitive neuroscience") and the level of complexity appropriate for your audience (e.g. "undergraduate level slides").
    prompt = (
        "You are an AI assistant helping a professor create speaker’s notes for a set of lecture slides on an academic topic.\n"
        "For each slide, please generate 2–3 paragraphs of clear, engaging, and explanatory speaker’s notes. Make sure the complexity of the content is appropriate for a group of undergraduate students and/or professionals. If possible, include examples, quotes, and elaborations on the content within each slide. Format using bullet points when appropriate."
        "Here are the slides:\n\n"
    )
    
    for i, slide_text in enumerate(slides_text_chunk, start=1):
        prompt += f"Slide {i}:\n{slide_text}\n\n"
    
    prompt += (
        "For each slide above, return the speaker’s notes labeled by slide number. "
    # You can also modify the prompt here, per-slide. Useful if you want to ensure a specific type of formatting/tone across all slides.
        "Use 2–3 paragraphs for each, in a conversational style.\n"
    )
    
    # Build the JSON payload
    payload = {
        "model": "gpt-4o",
        "messages": [
            {
                "role": "user",
                "content": prompt
            }
        ]
    }
    payload_str = json.dumps(payload)
    
    # Build and execute the curl command
    curl_command = [
        "curl",
        "https://api.openai.com/v1/chat/completions",
        "-H", "Content-Type: application/json",
        "-H", f"Authorization: Bearer {OPENAI_API_KEY}",
        "-d", payload_str
    ]
    
    try:
        result = subprocess.check_output(curl_command)
    except subprocess.CalledProcessError as e:
        error_output = e.output.decode("utf-8") if e.output else "No output"
        raise Exception(f"Curl command failed with return code {e.returncode}: {error_output}")
    
    result_str = result.decode("utf-8")
    
    try:
        response = json.loads(result_str)
    except json.JSONDecodeError as e:
        raise Exception(f"Failed to decode JSON from API response. Response text: {result_str}. Error: {e}")
    
    # Check for API errors
    if "error" in response:
        raise Exception(f"API returned an error: {response['error']}. Full response: {response}")
    
    if "choices" not in response:
        raise Exception(f"Unexpected API response format, 'choices' key missing. Full response: {response}")
    
    try:
        notes_text = response["choices"][0]["message"]["content"]
    except (IndexError, KeyError) as e:
        raise Exception(f"Error extracting notes text from response. Full response: {response}. Error: {e}")
    
    return notes_text

# 5) FUNCTION TO INSERT NOTES INTO THE PPTX
def insert_notes_to_pptx(original_pptx_path, notes_dict):
    """
    notes_dict maps 0-based slide indices to speaker notes text.
    We first ensure every slide has a notes slide (with a ' ' placeholder)
    so that python-pptx won't skip slides without existing notes.
    """
    prs = Presentation(original_pptx_path)

    # ——————————————————————————————
    # 0) Ensure every slide has a notes slide with at least one character
    for slide in prs.slides:
        # Accessing slide.notes_slide will create it if missing
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        if not tf.text:
            tf.text = " "  # single‐space placeholder
    # ——————————————————————————————

    # 1) Now overwrite with real notes where available
    for slide_index, slide in enumerate(prs.slides):
        if slide_index in notes_dict:
            notes_slide = slide.notes_slide  # now guaranteed to exist
            notes_slide.notes_text_frame.text = notes_dict[slide_index]

    return prs


# 6) MAIN SCRIPT TO PROCESS A FOLDER OF PPTX FILES
def process_pptx_folder(input_folder, output_folder, chunk_size=5):
    # Ensure output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for filename in os.listdir(input_folder):
        if filename.lower().endswith(".pptx"):
            pptx_path = os.path.join(input_folder, filename)
            slides_text = extract_slides_text(pptx_path)
            
            # Dictionary to hold notes for each slide
            all_notes = {}
            
            # Process slides in chunks to manage token limits
            for chunk_index, chunk in enumerate(chunk_list(slides_text, chunk_size)):
                try:
                    notes_response = get_speaker_notes(chunk)
                except Exception as e:
                    print(f"Error processing slides chunk starting at slide {chunk_index * chunk_size + 1} in file '{filename}': {e}")
                    continue
                
                # Use regex to extract slide numbers and corresponding notes
                # This regex captures: Slide <number>: <notes text> until the next "Slide <number>:" or end-of-text.
                pattern = re.compile(r"Slide\s*(\d+):\s*(.*?)(?=Slide\s*\d+:|$)", re.DOTALL)
                matches = pattern.findall(notes_response)
                if not matches:
                    print(f"No valid slide notes found in API response for chunk starting at slide {chunk_index * chunk_size + 1} in file '{filename}'. Response: {notes_response}")
                
                for slide_num_str, note_text in matches:
                    try:
                        # Calculate the actual slide index in the full PPTX
                        slide_index = int(slide_num_str) - 1 + (chunk_index * chunk_size)
                        all_notes[slide_index] = note_text.strip()
                    except Exception as e:
                        print(f"Error parsing slide number '{slide_num_str}' in file '{filename}': {e}")
            
            # Insert the speaker notes into the PPTX
            updated_pptx = insert_notes_to_pptx(pptx_path, all_notes)
            
            # Save the updated file
            output_path = os.path.join(output_folder, filename)
            updated_pptx.save(output_path)
            print(f"Processed and saved: {output_path}")

#Adjust the chunk_size variable below to change how many slides are processed per chunk. Lower values generally have higher quality outputs but also make a greater volume of API calls.
if __name__ == "__main__":
    input_folder = "C:/YOUR/INPUT/FOLDER/PATH/input_pptx"
    output_folder = "C:/YOUR/OUTPUT/FOLDER/PATH/output_pptx"
    process_pptx_folder(input_folder, output_folder, chunk_size=5)
